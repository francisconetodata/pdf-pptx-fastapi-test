from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, Dict, Any
import PyPDF2
from pptx import Presentation
import io
from datetime import datetime
import hashlib
import json  # <--- ADICIONADO
import re    # <--- ADICIONADO

app = FastAPI(title="Document Extractor API", version="1.0.0")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class DocumentResponse(BaseModel):
    success: bool
    filename: str
    file_type: str
    text_content: str
    metadata: Dict[str, Any]
    extracted_at: str
    file_hash: str

class TextInput(BaseModel):  # <--- ADICIONADO
    """Modelo para receber o texto bruto que contém JSON"""
    raw_text: str

def calculate_file_hash(content: bytes) -> str:
    """Calcula hash SHA256 do arquivo"""
    return hashlib.sha256(content).hexdigest()

def extract_pdf(file_content: bytes, filename: str) -> DocumentResponse:
    """Extrai texto e metadados de PDF"""
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        # Extrai texto de todas as páginas
        text_content = ""
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            text_content += f"\n--- Página {page_num + 1} ---\n{text}"
        
        # Extrai metadados
        metadata = {
            "num_pages": len(pdf_reader.pages),
            "pdf_metadata": {}
        }
        
        if pdf_reader.metadata:
            metadata["pdf_metadata"] = {
                "title": pdf_reader.metadata.get('/Title', ''),
                "author": pdf_reader.metadata.get('/Author', ''),
                "subject": pdf_reader.metadata.get('/Subject', ''),
                "creator": pdf_reader.metadata.get('/Creator', ''),
                "producer": pdf_reader.metadata.get('/Producer', ''),
                "creation_date": str(pdf_reader.metadata.get('/CreationDate', '')),
                "modification_date": str(pdf_reader.metadata.get('/ModDate', ''))
            }
        
        file_hash = calculate_file_hash(file_content)
        
        return DocumentResponse(
            success=True,
            filename=filename,
            file_type="pdf",
            text_content=text_content.strip(),
            metadata=metadata,
            extracted_at=datetime.utcnow().isoformat(),
            file_hash=file_hash
        )
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Erro ao processar PDF: {str(e)}")

def extract_pptx(file_content: bytes, filename: str) -> DocumentResponse:
    """Extrai texto e metadados de PPTX"""
    try:
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        # Extrai texto de todos os slides
        text_content = ""
        slide_count = 0
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += f"{shape.text}\n"
            
            text_content += slide_text
            slide_count += 1
        
        # Extrai metadados
        metadata = {
            "num_slides": slide_count,
            "slide_width": presentation.slide_width,
            "slide_height": presentation.slide_height,
            "pptx_metadata": {}
        }
        
        # Propriedades do core
        core_props = presentation.core_properties
        metadata["pptx_metadata"] = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "comments": core_props.comments or "",
            "category": core_props.category or "",
            "created": core_props.created.isoformat() if core_props.created else "",
            "modified": core_props.modified.isoformat() if core_props.modified else "",
            "last_modified_by": core_props.last_modified_by or ""
        }
        
        file_hash = calculate_file_hash(file_content)
        
        return DocumentResponse(
            success=True,
            filename=filename,
            file_type="pptx",
            text_content=text_content.strip(),
            metadata=metadata,
            extracted_at=datetime.utcnow().isoformat(),
            file_hash=file_hash
        )
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Erro ao processar PPTX: {str(e)}")

@app.get("/")
def read_root():
    return {
        "message": "Document Extractor API",
        "version": "1.0.0",
        "endpoints": {
            "/extract": "POST - Upload PDF ou PPTX para extração",
            "/transform-to-json": "POST - Extrai JSON de um bloco de texto", # <--- ATUALIZADO
            "/health": "GET - Status da API"
        }
    }

@app.get("/health")
def health_check():
    return {"status": "healthy", "timestamp": datetime.utcnow().isoformat()}

@app.post("/extract", response_model=DocumentResponse)
async def extract_document(file: UploadFile = File(...)):
    """
    Extrai texto e metadados de arquivos PDF ou PPTX
    """
    # Valida tipo de arquivo
    file_extension = file.filename.lower().split('.')[-1]
    
    if file_extension not in ['pdf', 'pptx']:
        raise HTTPException(
            status_code=400, 
            detail="Tipo de arquivo não suportado. Use PDF ou PPTX."
        )
    
    # Lê conteúdo do arquivo
    file_content = await file.read()
    
    if len(file_content) == 0:
        raise HTTPException(status_code=400, detail="Arquivo vazio")
    
    # Processa conforme tipo
    if file_extension == 'pdf':
        return extract_pdf(file_content, file.filename)
    elif file_extension == 'pptx':
        return extract_pptx(file_content, file.filename)

# --- INÍCIO DA NOVA ROTA ---
@app.post("/transform-to-json", response_model=Dict[str, Any])
async def transform_text_to_json(payload: TextInput):
    """
    Recebe um texto que contém uma estrutura JSON (potencialmente
    dentro de um array e/ou bloco de markdown) e extrai o JSON interno.
    
    Espera um payload como: {"raw_text": "[{\\"text\\":\\"```json...```\\"}]"}
    """
    try:
        # 1. Tenta decodificar o texto bruto (que esperamos ser um JSON de lista)
        data_list = json.loads(payload.raw_text)
        
        # 2. Validações da estrutura esperada
        if not isinstance(data_list, list) or len(data_list) == 0:
            raise ValueError("Input 'raw_text' não é uma lista JSON válida ou está vazia.")
        
        first_item = data_list[0]
        if not isinstance(first_item, dict) or "text" not in first_item:
            raise ValueError("O primeiro item da lista JSON não contém a chave 'text'.")
            
        text_with_markdown = first_item["text"]
        
        # 3. Extrai o JSON de dentro do bloco de markdown ```json ... ```
        # re.DOTALL faz o '.' corresponder a quebras de linha
        match = re.search(r"```json\s*(.*?)\s*```", text_with_markdown, re.DOTALL | re.IGNORECASE)
        
        json_string = ""
        if match:
            json_string = match.group(1) # Pega o que está dentro do bloco
        else:
            # Se não achar o bloco, levanta um erro
            raise ValueError("Não foi possível encontrar o bloco ```json ... ``` dentro do campo 'text'.")
        
        # 4. Decodifica o JSON extraído
        final_json_data = json.loads(json_string)
        
        # 5. Retorna o JSON final
        return final_json_data

    except json.JSONDecodeError as e:
        # Erro ao tentar decodificar 'raw_text' ou o 'json_string' extraído
        raise HTTPException(status_code=400, detail=f"Erro ao decodificar JSON: {str(e)}")
    except ValueError as e:
        # Erro das nossas validações
        raise HTTPException(status_code=400, detail=f"Erro na estrutura do input: {str(e)}")
    except Exception as e:
        # Outros erros
        raise HTTPException(status_code=500, detail=f"Erro inesperado no processamento: {str(e)}")
# --- FIM DA NOVA ROTA ---

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)